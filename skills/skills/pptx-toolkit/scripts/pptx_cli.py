#!/usr/bin/env python3
"""pptx-toolkit CLI — read/outline/notes/to-md subcommands.

MVP scope (read-only):
- read:    slide text + notes → markdown
- outline: slide titles (numbered)
- notes:   speaker notes only
- to-md:   markdown aggregator (same as read --output)

Uses python-pptx 1.0.2 from ~/.anchor/env/.venv.
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import List, Optional

try:
    from pptx import Presentation
    from pptx.slide import Slide
except ImportError:
    print("[pptx-toolkit] python-pptx 미설치. ~/.anchor/env/.venv 확인", file=sys.stderr)
    sys.exit(2)


# ──────────────────────────────────────────────────────
# Slide extraction
# ──────────────────────────────────────────────────────
def load_pres(path: Path) -> Presentation:
    try:
        return Presentation(str(path))
    except Exception as e:
        print(f"[pptx-toolkit] python-pptx load 실패: {e}", file=sys.stderr)
        print("  .ppt (97-2003) 는 지원 안 됨. pandoc/soffice 변환 필요.", file=sys.stderr)
        sys.exit(2)


def slide_title(slide: Slide) -> str:
    """Best-effort extraction of a slide title."""
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        t = slide.shapes.title.text_frame.text.strip()
        if t:
            return t
    # Fallback: first text_frame found
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                # Take first line
                return t.splitlines()[0][:120]
    return "(untitled)"


def slide_body(slide: Slide) -> List[str]:
    """Extract all non-title text blocks as flat list of lines."""
    lines: List[str] = []
    title_shape = slide.shapes.title
    title_id = title_shape.shape_id if title_shape is not None else None
    for shape in slide.shapes:
        if title_id is not None and shape.shape_id == title_id:
            continue
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = "".join(run.text for run in para.runs).strip()
            if text:
                # Indent based on paragraph level
                indent = "  " * max(0, para.level)
                lines.append(f"{indent}- {text}")
    return lines


def slide_notes(slide: Slide) -> Optional[str]:
    if not slide.has_notes_slide:
        return None
    frame = slide.notes_slide.notes_text_frame
    text = frame.text.strip() if frame else ""
    return text or None


def slide_shape_summary(slide: Slide) -> str:
    """Count shapes by type."""
    counts: dict[str, int] = {}
    for shape in slide.shapes:
        name = shape.shape_type.__str__() if shape.shape_type else "UNKNOWN"
        counts[name] = counts.get(name, 0) + 1
    if not counts:
        return ""
    return ", ".join(f"{k}×{v}" for k, v in sorted(counts.items()))


# ──────────────────────────────────────────────────────
# Formatters
# ──────────────────────────────────────────────────────
def format_slide_md(idx: int, slide: Slide, include_notes: bool = True) -> str:
    out = [f"## Slide {idx}: {_md_escape(slide_title(slide))}", ""]
    body = slide_body(slide)
    if body:
        out.extend(body)
        out.append("")
    if include_notes:
        notes = slide_notes(slide)
        if notes:
            out.append("**발표자 노트:**")
            out.append("")
            for line in notes.splitlines():
                out.append(f"> {line}")
            out.append("")
    return "\n".join(out)


def _md_escape(s: str) -> str:
    return s.replace("|", r"\|").strip()


# ──────────────────────────────────────────────────────
# Subcommands
# ──────────────────────────────────────────────────────
def cmd_read(args: argparse.Namespace) -> int:
    pres = load_pres(Path(args.file))
    total = len(pres.slides)
    parts = [f"# {Path(args.file).name}", f"- Slides: {total}", ""]
    if args.slide is not None:
        if not (1 <= args.slide <= total):
            print(f"[pptx-toolkit] 슬라이드 {args.slide} 없음 (1..{total})", file=sys.stderr)
            return 3
        parts.append(format_slide_md(args.slide, pres.slides[args.slide - 1]))
    else:
        for i, slide in enumerate(pres.slides, start=1):
            parts.append(format_slide_md(i, slide))
    _emit("\n".join(parts), args.output)
    return 0


def cmd_outline(args: argparse.Namespace) -> int:
    pres = load_pres(Path(args.file))
    total = len(pres.slides)
    out = [f"# {Path(args.file).name} — Outline ({total} slides)", ""]
    for i, slide in enumerate(pres.slides, start=1):
        title = _md_escape(slide_title(slide))
        out.append(f"{i:3}. {title}")
    _emit("\n".join(out) + "\n", args.output)
    return 0


def cmd_notes(args: argparse.Namespace) -> int:
    pres = load_pres(Path(args.file))
    total = len(pres.slides)
    parts = [f"# {Path(args.file).name} — Speaker Notes", ""]
    if args.slide is not None:
        if not (1 <= args.slide <= total):
            print(f"[pptx-toolkit] 슬라이드 {args.slide} 없음", file=sys.stderr)
            return 3
        slides = [(args.slide, pres.slides[args.slide - 1])]
    else:
        slides = list(enumerate(pres.slides, start=1))
    for i, slide in slides:
        notes = slide_notes(slide)
        parts.append(f"## Slide {i}: {_md_escape(slide_title(slide))}")
        parts.append("")
        if notes:
            for line in notes.splitlines():
                parts.append(f"> {line}")
        else:
            parts.append("_(no notes)_")
        parts.append("")
    _emit("\n".join(parts), args.output)
    return 0


def cmd_to_md(args: argparse.Namespace) -> int:
    # Same as `read` but always writes to output file if -o given
    return cmd_read(args)


def _emit(content: str, output: Optional[str]) -> None:
    if output:
        Path(output).write_text(content, encoding="utf-8")
        print(f"wrote {output}", file=sys.stderr)
    else:
        sys.stdout.write(content)


# ──────────────────────────────────────────────────────
# Argparse
# ──────────────────────────────────────────────────────
def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(prog="pptx", description="pptx-toolkit — PowerPoint read-only CLI")
    sub = p.add_subparsers(dest="command", required=True)

    def _add_file(sp):
        sp.add_argument("file", help=".pptx file path")
        sp.add_argument("-o", "--output", help="output file (default stdout)")

    sp_read = sub.add_parser("read", help="slide text + notes → markdown")
    _add_file(sp_read)
    sp_read.add_argument("--slide", type=int, help="specific slide number (1-based)")
    sp_read.set_defaults(func=cmd_read)

    sp_outline = sub.add_parser("outline", help="slide titles (numbered)")
    _add_file(sp_outline)
    sp_outline.set_defaults(func=cmd_outline)

    sp_notes = sub.add_parser("notes", help="speaker notes only")
    _add_file(sp_notes)
    sp_notes.add_argument("--slide", type=int, help="specific slide number (1-based)")
    sp_notes.set_defaults(func=cmd_notes)

    sp_md = sub.add_parser("to-md", help="markdown export (alias for read)")
    _add_file(sp_md)
    sp_md.add_argument("--slide", type=int, help="specific slide number (1-based)")
    sp_md.set_defaults(func=cmd_to_md)

    return p


def main() -> int:
    parser = build_parser()
    args = parser.parse_args()
    if not Path(args.file).exists():
        print(f"[pptx-toolkit] 파일 없음: {args.file}", file=sys.stderr)
        return 1
    return args.func(args)


if __name__ == "__main__":
    sys.exit(main())
